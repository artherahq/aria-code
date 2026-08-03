"""report_exporters.py — Markdown → Word (.docx) / PowerPoint (.pptx) export.

Complements markdown_pdf.py (Markdown → PDF): same input, different output
formats, for callers who need an editable document instead of print output.

Parsing approach: run the input through the `markdown` library (already a
core dependency) to get real HTML, then walk that with BeautifulSoup (already
a "files"-extra dependency, used elsewhere for document parsing) to pull out
a flat sequence of (heading-level | paragraph | bullet) blocks. This only
covers the markdown subset aria-code's own report generators actually emit
(headings, paragraphs, bullet/numbered lists, bold/italic inline) — it is not
a full CommonMark-to-OOXML converter, and unsupported constructs (tables,
images, nested blockquotes) are dropped rather than mis-rendered.

Public API:
    markdown_to_docx(md_text_or_path, docx_path, title=None) -> Optional[Path]
    markdown_to_pptx(md_text_or_path, pptx_path, title=None) -> Optional[Path]
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

logger = logging.getLogger(__name__)


def _read_md(md_text_or_path: Union[str, Path]) -> str:
    if isinstance(md_text_or_path, Path) or (
        isinstance(md_text_or_path, str) and len(md_text_or_path) < 260 and Path(md_text_or_path).exists()
    ):
        return Path(md_text_or_path).read_text(encoding="utf-8")
    return str(md_text_or_path)


def _parse_blocks(md_text: str) -> List[Dict[str, Any]]:
    """Return a flat list of {"type": "heading"|"para"|"bullet", "level": int, "text": str}."""
    import markdown as _markdown
    from bs4 import BeautifulSoup

    html = _markdown.markdown(md_text, extensions=["extra"])
    soup = BeautifulSoup(html, "html.parser")

    blocks: List[Dict[str, Any]] = []
    for el in soup.find_all(["h1", "h2", "h3", "h4", "p", "ul", "ol"], recursive=False):
        if el.name in ("h1", "h2", "h3", "h4"):
            text = el.get_text(strip=True)
            if text:
                blocks.append({"type": "heading", "level": int(el.name[1]), "text": text})
        elif el.name == "p":
            text = el.get_text(strip=True)
            if text:
                blocks.append({"type": "para", "text": text})
        elif el.name in ("ul", "ol"):
            for li in el.find_all("li", recursive=False):
                text = li.get_text(strip=True)
                if text:
                    blocks.append({"type": "bullet", "text": text})
    return blocks


def markdown_to_docx(
    md_text_or_path: Union[str, Path],
    docx_path: Union[str, Path],
    title: Optional[str] = None,
) -> Optional[Path]:
    """Render Markdown to a Word document. Returns the output path, or None on failure."""
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        logger.warning("python-docx not installed — pip install python-docx")
        return None

    md_text = _read_md(md_text_or_path)
    blocks = _parse_blocks(md_text)
    docx_path = Path(docx_path)

    doc = Document()
    if title:
        doc.add_heading(title, level=0)

    heading_styles = {1: "Heading 1", 2: "Heading 2", 3: "Heading 3", 4: "Heading 4"}
    for block in blocks:
        if block["type"] == "heading":
            doc.add_heading(block["text"], level=min(block["level"], 4))
        elif block["type"] == "bullet":
            doc.add_paragraph(block["text"], style="List Bullet")
        else:
            p = doc.add_paragraph(block["text"])
            p.paragraph_format.space_after = Pt(8)

    docx_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(docx_path))
    return docx_path


def markdown_to_pptx(
    md_text_or_path: Union[str, Path],
    pptx_path: Union[str, Path],
    title: Optional[str] = None,
) -> Optional[Path]:
    """Render Markdown to a slide deck: each heading starts a new slide, with
    following paragraphs/bullets as its body. Returns the output path, or None on failure."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        logger.warning("python-pptx not installed — pip install python-pptx")
        return None

    md_text = _read_md(md_text_or_path)
    blocks = _parse_blocks(md_text)
    pptx_path = Path(pptx_path)

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]

    if title:
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title

    current_slide = None
    current_body = None
    for block in blocks:
        if block["type"] == "heading":
            current_slide = prs.slides.add_slide(body_layout)
            current_slide.shapes.title.text = block["text"]
            current_body = current_slide.placeholders[1].text_frame
            current_body.clear()
            first_line = True
        else:
            if current_slide is None:
                # Content before the first heading — give it its own slide.
                current_slide = prs.slides.add_slide(body_layout)
                current_slide.shapes.title.text = title or "Summary"
                current_body = current_slide.placeholders[1].text_frame
                current_body.clear()
                first_line = True
            para = current_body.paragraphs[0] if first_line else current_body.add_paragraph()
            para.text = block["text"]
            para.level = 1 if block["type"] == "bullet" else 0
            para.font.size = Pt(16)
            first_line = False

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    return pptx_path
